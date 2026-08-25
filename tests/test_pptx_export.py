"""VISION.md 対応: 機能別ルール(PPTX生成)のレイアウト計算。"""

from __future__ import annotations

from PIL import Image
from pptx import Presentation

import pptx_export


def test_grid_shape_for_small_counts():
    assert pptx_export._grid_shape(1) == (1, 1)
    assert pptx_export._grid_shape(3) == (3, 1)


def test_grid_shape_for_medium_counts():
    assert pptx_export._grid_shape(4) == (4, 1)
    assert pptx_export._grid_shape(8) == (4, 2)


def test_grid_shape_for_large_counts():
    assert pptx_export._grid_shape(9) == (5, 2)
    assert pptx_export._grid_shape(20) == (5, 4)


def test_export_places_one_picture_and_one_label_per_logo(tmp_path):
    entries = []
    for i, name in enumerate(["あ社", "い社", "う社"]):
        p = tmp_path / f"{name}.png"
        Image.new("RGBA", (300, 100), (10 + i, 20, 30, 255)).save(p)
        entries.append((name, p))

    out = pptx_export.export(entries, tmp_path / "logos.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    # ロゴ1件につき「画像1つ + ラベル1つ」の2図形
    assert len(prs.slides[0].shapes) == len(entries) * 2


def test_export_with_no_entries_still_produces_a_valid_file(tmp_path):
    out = pptx_export.export([], tmp_path / "empty.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert len(prs.slides[0].shapes) == 0
