"""全ロゴを1スライドにグリッド配置した pptx を出力する。

クリップボード連携が使えない環境でのフォールバックも兼ねる。
PowerPoint で開き、必要なロゴを選んでコピーすればよい。
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
LABEL_H = Inches(0.3)
GAP = Inches(0.25)


def _grid_shape(count: int) -> tuple[int, int]:
    """件数から列数・行数を決める。10社以下を想定し、横長に並べる。"""
    if count <= 3:
        return count, 1
    if count <= 8:
        return 4, (count + 3) // 4
    return 5, (count + 4) // 5


def export(entries: list[tuple[str, Path]], out_path: Path) -> Path:
    """entries は (企業名, 透過PNGのパス) の並び。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙レイアウト

    if not entries:
        return _save(prs, out_path)

    cols, rows = _grid_shape(len(entries))
    cell_w = (SLIDE_W - MARGIN * 2 - GAP * (cols - 1)) / cols
    cell_h = (SLIDE_H - MARGIN * 2 - GAP * (rows - 1)) / rows
    image_h = cell_h - LABEL_H

    for i, (name, png_path) in enumerate(entries):
        col, row = i % cols, i // cols
        cell_x = MARGIN + (cell_w + GAP) * col
        cell_y = MARGIN + (cell_h + GAP) * row

        with Image.open(png_path) as im:
            iw, ih = im.size
        # セルに収まるよう縦横比を保って縮める
        scale = min(cell_w / iw, image_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(
            str(png_path),
            Emu(int(cell_x + (cell_w - w) / 2)),
            Emu(int(cell_y + (image_h - h) / 2)),
            width=Emu(w),
            height=Emu(h),
        )

        box = slide.shapes.add_textbox(
            Emu(int(cell_x)), Emu(int(cell_y + image_h)), Emu(int(cell_w)), Emu(int(LABEL_H))
        )
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    return _save(prs, out_path)


def _save(prs: Presentation, out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
